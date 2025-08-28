
# import streamlit as st
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# import io, json, re, requests
# from typing import List, Dict, Any, Optional

# st.set_page_config(page_title="Your Text, Your Style", page_icon="📊", layout="wide")

# # -------------------------
# # Security: do not log keys, avoid storing secrets
# # -------------------------
# def red(text):  # tiny helper
#     st.markdown(f"<span style='color:#d00'>{text}</span>", unsafe_allow_html=True)

# # -------------------------
# # LLM Calls
# # -------------------------
# def call_llm(provider: str, api_key: str, model: str, system_prompt: str, user_prompt: str, temperature: float = 0.2) -> str:
#     """Call the chosen provider; return text. Minimal error surfacing; no key logging."""
#     try:
#         if provider == "OpenAI":
#             url = "https://api.openai.com/v1/chat/completions"
#             headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
#             data = {
#                 "model": model or "gpt-4o-mini",
#                 "messages": [
#                     {"role":"system","content": system_prompt},
#                     {"role":"user","content": user_prompt},
#                 ],
#                 "temperature": temperature,
#             }
#             r = requests.post(url, headers=headers, json=data, timeout=60)
#             r.raise_for_status()
#             return r.json()["choices"][0]["message"]["content"]

#         elif provider == "Anthropic":
#             url = "https://api.anthropic.com/v1/messages"
#             headers = {"x-api-key": api_key, "anthropic-version":"2023-06-01","content-type":"application/json"}
#             data = {
#                 "model": model or "claude-3-haiku-20240307",
#                 "max_tokens": 2000,
#                 "temperature": temperature,
#                 "system": system_prompt,
#                 "messages": [{"role":"user","content": user_prompt}],
#             }
#             r = requests.post(url, headers=headers, json=data, timeout=60)
#             r.raise_for_status()
#             out = r.json()
#             return "".join([b.get("text","") for b in out.get("content", [])])

#         elif provider == "Gemini":
#             use_model = model or "gemini-1.5-flash"
#             url = f"https://generativelanguage.googleapis.com/v1beta/models/{use_model}:generateContent?key={api_key}"
#             data = {
#                 "contents":[
#                     {"parts":[{"text": system_prompt}]},
#                     {"parts":[{"text": user_prompt}]},
#                 ],
#                 "generationConfig":{"temperature": temperature},
#             }
#             r = requests.post(url, json=data, timeout=60)
#             r.raise_for_status()
#             out = r.json()
#             cands = out.get("candidates", [])
#             if not cands:
#                 return ""
#             parts = cands[0].get("content", {}).get("parts", [])
#             return "".join(p.get("text","") for p in parts)

#         return ""
#     except Exception as e:
#         st.error("LLM call failed. Check provider/model/key or your network.")
#         return ""

# def robust_json_extract(text: str) -> Optional[Dict[str, Any]]:
#     """Try to parse a JSON object from LLM text output."""
#     if not text:
#         return None
#     try:
#         return json.loads(text)
#     except Exception:
#         pass
#     m = re.search(r"\{.*\}", text, flags=re.S)
#     if m:
#         try:
#             return json.loads(m.group(0))
#         except Exception:
#             return None
#     return None

# LLM_SYSTEM = """You are a slide architect. Transform the user's text into a clean, concise deck.
# Return ONLY valid JSON with this schema:
# {
#   "slides": [
#     { "title": "Concise title",
#       "bullets": ["bullet 1", "bullet 2", "..."],
#       "notes": "optional speaker notes" }
#   ]
# }
# Choose a reasonable number of slides (6–20) based on input length and guidance.
# Prefer 3–6 bullets per slide. Use plain text only (no markdown, no numbering). Keep titles concise.
# """

# def naive_chunk_to_slides(text: str, max_chars: int = 700) -> Dict[str, Any]:
#     """Fallback splitter: split text into slides using headings then paragraphs."""
#     slides = []
#     text = text.strip()
#     if not text:
#         return {"slides": []}
#     sections = re.split(r"\n(?=# )", text)
#     if len(sections) == 1:
#         paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
#         current = []
#         acc = 0
#         for p in paras:
#             if acc + len(p) > max_chars and current:
#                 title = current[0][:70]
#                 bullets = current[1:5] if len(current) > 1 else []
#                 slides.append({"title": title, "bullets": bullets, "notes": ""})
#                 current, acc = [], 0
#             current.append(p)
#             acc += len(p)
#         if current:
#             title = current[0][:70]
#             bullets = current[1:5] if len(current) > 1 else []
#             slides.append({"title": title, "bullets": bullets, "notes": ""})
#     else:
#         for sec in sections:
#             lines = [l.strip() for l in sec.splitlines() if l.strip()]
#             title = lines[0].lstrip("# ").strip() if lines else "Slide"
#             bullets = [l for l in lines[1:6]]
#             slides.append({"title": title, "bullets": bullets, "notes": ""})
#     return {"slides": slides[:25]}

# def infer_slides_with_llm(provider, api_key, model, text, guidance, gen_notes=True):
#     user_prompt = f"Guidance: {guidance or 'none'}\n\nSource text:\n{text}"
#     out = call_llm(provider, api_key, model, LLM_SYSTEM, user_prompt, temperature=0.2)
#     data = robust_json_extract(out) if out else None
#     if not data or "slides" not in data:
#         data = naive_chunk_to_slides(text)
#     if not gen_notes:
#         for s in data.get("slides", []):
#             s["notes"] = ""
#     return data

# # -------------------------
# # Template helpers
# # -------------------------
# def delete_all_slides(prs: Presentation):
#     ids = list(prs.slides._sldIdLst)
#     for sld in ids:
#         prs.part.drop_rel(sld.rId)
#         prs.slides._sldIdLst.remove(sld)

# def collect_template_images(prs: Presentation):
#     seen = {}
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#                 img = shape.image
#                 seen[img.filename] = img.blob
#     return list(seen.items())

# def choose_layout(prs: Presentation, want_content=True):
#     preferred = "Title and Content" if want_content else "Title Only"
#     for layout in prs.slide_layouts:
#         if layout.name == preferred:
#             return layout
#     # fallbacks
#     for layout in prs.slide_layouts:
#         if "Title" in layout.name and "Content" in layout.name:
#             return layout
#     for layout in prs.slide_layouts:
#         if "Title" in layout.name:
#             return layout
#     return prs.slide_layouts[0]

# def fill_placeholders(slide, title: str, bullets: List[str]):
#     title_set, content_set = False, False
#     for shape in slide.placeholders:
#         phf = getattr(shape, "placeholder_format", None)
#         if not phf:
#             continue
#         if shape.has_text_frame:
#             # Title placeholders have type 1
#             if not title_set and (phf.type == 1 or "title" in shape.name.lower()):
#                 shape.text_frame.clear()
#                 shape.text_frame.text = title
#                 title_set = True
#             elif not content_set and bullets and (phf.type == 2 or "content" in shape.name.lower()):
#                 tf = shape.text_frame
#                 tf.clear()
#                 tf.text = bullets[0]
#                 for b in bullets[1:]:
#                     p = tf.add_paragraph()
#                     p.text = b
#                     p.level = 0
#                 content_set = True
#     if bullets and not content_set:
#         # Add a text box if no content placeholder was found
#         tb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
#         tf = tb.text_frame
#         tf.text = bullets[0]
#         for b in bullets[1:]:
#             p = tf.add_paragraph()
#             p.text = b
#             p.level = 0

# def add_notes(slide, notes: str):
#     if notes:
#         ns = slide.notes_slide
#         ns.notes_text_frame.text = notes

# def build_presentation(template_bytes: bytes, slides_data: Dict[str, Any]) -> bytes:
#     base = Presentation(io.BytesIO(template_bytes))
#     # Re-open template to collect images without altering the working deck
#     template_for_images = Presentation(io.BytesIO(template_bytes))
#     images = collect_template_images(template_for_images)
#     delete_all_slides(base)

#     for i, s in enumerate(slides_data.get("slides", [])):
#         title = s.get("title", "Slide")
#         bullets = [b.strip() for b in (s.get("bullets") or []) if b and b.strip()]
#         notes = s.get("notes", "")
#         layout = choose_layout(base, want_content=bool(bullets))
#         slide = base.slides.add_slide(layout)
#         fill_placeholders(slide, title, bullets)
#         add_notes(slide, notes)

#         # Opportunistic image reuse: add one of the template images in a corner (if any)
#         if images:
#             try:
#                 fname, blob = images[i % len(images)]
#                 slide.shapes.add_picture(io.BytesIO(blob), Inches(8.2), Inches(4.6), width=Inches(1.5))
#             except Exception:
#                 pass

#     out = io.BytesIO()
#     base.save(out)
#     return out.getvalue()

# # -------------------------
# # UI
# # -------------------------
# st.title("📊 Your Text, Your Style – Auto-Generate a Presentation")
# st.caption("Paste text or markdown, upload a PowerPoint template, and get a styled .pptx. Your API key is never stored or logged.")

# with st.form("form"):
#     c1, c2 = st.columns([2,1])
#     with c1:
#         text = st.text_area("Input text or markdown", height=260, placeholder="Paste chapters, notes, or long-form prose...")
#         guidance = st.text_input("Optional guidance (tone/use case)", placeholder="e.g., 'turn into an investor pitch deck'")
#     with c2:
#         provider = st.selectbox("LLM Provider", ["OpenAI","Anthropic","Gemini"])
#         model = st.text_input("Model (optional)", placeholder={
#             "OpenAI":"gpt-4o-mini / gpt-4o",
#             "Anthropic":"claude-3-haiku-20240307",
#             "Gemini":"gemini-1.5-flash"
#         }[provider])
#         api_key = st.text_input(f"{provider} API Key", type="password", help="Used only for this request, not stored.")
#         template = st.file_uploader("Upload .pptx or .potx", type=["pptx","potx"])
#         gen_notes = st.checkbox("Auto-generate speaker notes", value=True)
#     submitted = st.form_submit_button("Generate Presentation")

# if submitted:
#     if not text or not template or not api_key:
#         st.warning("Please provide: input text, template file, and an API key.")
#         st.stop()
#     with st.spinner("Analyzing text and mapping to slides..."):
#         slides = infer_slides_with_llm(provider, api_key, model, text, guidance, gen_notes=gen_notes)
#     with st.spinner("Applying template styles & building PPTX..."):
#         pptx_bytes = build_presentation(template.read(), slides)
#     st.success("Done! Download your presentation:")
#     st.download_button("⬇️ Download .pptx", data=pptx_bytes, file_name="generated.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# st.markdown("---")
# with st.expander("Privacy & notes"):
#     st.write("""
# - Your API key is only used in-memory to call your chosen LLM and is never stored or logged by this app.
# - No images are generated by AI. The app can only reuse images found inside the uploaded template/presentation.
# - Layout fidelity is best-effort; exact placeholder mapping varies by template.
# """)


# import streamlit as st
# import requests
# import io
# from pptx import Presentation
# from pptx.util import Inches, Pt
# import time

# # -----------------------------
# # LLM Call Helper
# # -----------------------------
# def generate_completion_with_provider(provider, api_key, model, messages, max_retries=2, timeout=30):
#     if provider == "openai":
#         base = "https://api.openai.com/v1"
#         endpoint = f"{base}/chat/completions"
#         headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
#         payload = {"model": model, "messages": messages}

#     elif provider == "aipipe":
#         base = "https://aipipe.org/openai/v1"
#         endpoint = f"{base}/chat/completions"
#         headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
#         payload = {"model": model, "messages": messages}

#     elif provider == "aipipe_openrouter":
#         base = "https://aipipe.org/openrouter/v1"
#         endpoint = f"{base}/chat/completions"
#         headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
#         payload = {"model": model, "messages": messages}

#     elif provider == "anthropic":
#         base = "https://api.anthropic.com/v1/messages"
#         headers = {
#             "x-api-key": api_key,
#             "Content-Type": "application/json",
#             "anthropic-version": "2023-06-01",
#         }
#         payload = {"model": model, "messages": messages, "max_tokens": 1000}
#         endpoint = base

#     elif provider == "gemini":
#         base = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
#         headers = {"Content-Type": "application/json"}
#         # Gemini uses contents instead of messages
#         contents = [{"parts": [{"text": m["content"]}]} for m in messages if m["role"] == "user"]
#         payload = {"contents": contents}
#         endpoint = base

#     else:
#         raise ValueError("Unsupported provider")

#     attempt = 0
#     while True:
#         try:
#             r = requests.post(endpoint, json=payload, headers=headers, timeout=timeout)
#             r.raise_for_status()
#             return r.json()
#         except requests.RequestException as e:
#             attempt += 1
#             if attempt > max_retries:
#                 raise
#             time.sleep(1.5 ** attempt)

# # -----------------------------
# # Slide Builder
# # -----------------------------
# def build_presentation(template_file, slide_plan):
#     prs = Presentation(template_file)

#     for slide_data in slide_plan:
#         layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
#         slide = prs.slides.add_slide(layout)

#         if slide.shapes.title:
#             slide.shapes.title.text = slide_data.get("title", "")

#         body_shapes = [shape for shape in slide.placeholders if shape.is_placeholder and shape.placeholder_format.idx == 1]
#         if body_shapes:
#             tf = body_shapes[0].text_frame
#             tf.clear()
#             for bullet in slide_data.get("bullets", []):
#                 p = tf.add_paragraph()
#                 p.text = bullet
#                 p.level = 0

#     return prs

# # -----------------------------
# # Streamlit UI
# # -----------------------------
# st.title("📊 Your Text → PowerPoint Generator")

# provider = st.selectbox("LLM Provider", ["openai", "aipipe", "aipipe_openrouter", "anthropic", "gemini"])
# api_key = st.text_input("Paste your API key (never stored)", type="password")
# st.caption("👉 If using AIPipe, get a token at https://aipipe.org/login")

# uploaded_template = st.file_uploader("Upload a PowerPoint template (.pptx or .potx)", type=["pptx", "potx"])
# text_input = st.text_area("Paste your text / markdown / prose here:", height=300)
# guidance = st.text_input("Optional guidance (e.g. 'investor pitch deck')")

# if st.button("Generate Presentation"):
#     if not api_key:
#         st.error("Please paste your API key.")
#     elif not uploaded_template:
#         st.error("Please upload a template file.")
#     elif not text_input.strip():
#         st.error("Please paste some text to generate slides.")
#     else:
#         with st.spinner("Calling LLM and building slides..."):
#             try:
#                 messages = [
#                     {"role": "system", "content": "You are an assistant that turns long text into a JSON slide plan."},
#                     {"role": "user", "content": f"Input text:\n{text_input}\n\nGuidance:{guidance}\n\nOutput format: JSON array of slides, each with 'title' and 'bullets' (list of short strings)."}
#                 ]

#                 resp = generate_completion_with_provider(provider, api_key, model="gpt-4o-mini", messages=messages)

#                 # Different providers nest responses differently
#                 if provider in ["openai", "aipipe", "aipipe_openrouter"]:
#                     content = resp["choices"][0]["message"]["content"]
#                 elif provider == "anthropic":
#                     content = resp["content"][0]["text"]
#                 elif provider == "gemini":
#                     content = resp["candidates"][0]["content"]["parts"][0]["text"]
#                 else:
#                     content = "[]"

#                 import json
#                 try:
#                     slide_plan = json.loads(content)
#                 except:
#                     st.warning("LLM did not return valid JSON, falling back to one-slide output.")
#                     slide_plan = [{"title": "Summary", "bullets": [text_input[:200] + "..."]}]

#                 prs = build_presentation(uploaded_template, slide_plan)
#                 output = io.BytesIO()
#                 prs.save(output)
#                 st.success("✅ Presentation ready!")
#                 st.download_button("Download PPTX", data=output.getvalue(), file_name="generated.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

#             finally:
#                 # clear key
#                 del api_key



import streamlit as st
import requests
import io
from pptx import Presentation
from pptx.util import Inches, Pt
import time
import json
import re

# -----------------------------
# LLM Call Helper
# -----------------------------
def generate_completion_with_provider(provider, api_key, model, messages, max_retries=2, timeout=30):
    if provider == "openai":
        base = "https://api.openai.com/v1"
        endpoint = f"{base}/chat/completions"
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        payload = {"model": model, "messages": messages}

    elif provider == "aipipe":
        base = "https://aipipe.org/openai/v1"
        endpoint = f"{base}/chat/completions"
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        payload = {"model": model, "messages": messages}

    elif provider == "aipipe_openrouter":
        base = "https://aipipe.org/openrouter/v1"
        endpoint = f"{base}/chat/completions"
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        payload = {"model": model, "messages": messages}

    elif provider == "anthropic":
        base = "https://api.anthropic.com/v1/messages"
        headers = {
            "x-api-key": api_key,
            "Content-Type": "application/json",
            "anthropic-version": "2023-06-01",
        }
        payload = {"model": model, "messages": messages, "max_tokens": 1000}
        endpoint = base

    elif provider == "gemini":
        base = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
        headers = {"Content-Type": "application/json"}
        # Gemini uses contents instead of messages
        contents = [{"parts": [{"text": m["content"]}]} for m in messages if m["role"] == "user"]
        payload = {"contents": contents}
        endpoint = base

    else:
        raise ValueError("Unsupported provider")

    attempt = 0
    while True:
        try:
            r = requests.post(endpoint, json=payload, headers=headers, timeout=timeout)
            r.raise_for_status()
            return r.json()
        except requests.RequestException as e:
            attempt += 1
            if attempt > max_retries:
                raise
            time.sleep(1.5 ** attempt)

# -----------------------------
# JSON Repair Helper
# -----------------------------
def safe_json_parse(content, fallback):
    try:
        return json.loads(content)
    except:
        # Try to extract JSON array from text
        match = re.search(r'\[.*\]', content, re.S)
        if match:
            try:
                return json.loads(match.group(0))
            except:
                pass
        return fallback

# -----------------------------
# Slide Builder
# -----------------------------
def build_presentation(template_file, slide_plan):
    prs = Presentation(template_file)

    for slide_data in slide_plan:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")

        body_shapes = [shape for shape in slide.placeholders if shape.is_placeholder and shape.placeholder_format.idx == 1]
        if body_shapes:
            tf = body_shapes[0].text_frame
            tf.clear()
            for bullet in slide_data.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    return prs

# -----------------------------
# Streamlit UI
# -----------------------------
st.title("📊 Your Text → PowerPoint Generator")

provider = st.selectbox("LLM Provider", ["openai", "aipipe", "aipipe_openrouter", "anthropic", "gemini"])
api_key = st.text_input("Paste your API key (never stored)", type="password")
st.caption("👉 If using AIPipe, get a token at https://aipipe.org/login")

uploaded_template = st.file_uploader("Upload a PowerPoint template (.pptx or .potx)", type=["pptx", "potx"])
text_input = st.text_area("Paste your text / markdown / prose here:", height=300)
guidance = st.text_input("Optional guidance (e.g. 'investor pitch deck')")

if st.button("Generate Presentation"):
    if not api_key:
        st.error("Please paste your API key.")
    elif not uploaded_template:
        st.error("Please upload a template file.")
    elif not text_input.strip():
        st.error("Please paste some text to generate slides.")
    else:
        with st.spinner("Calling LLM and building slides..."):
            try:
                messages = [
                    {"role": "system", "content": """You are an assistant that ONLY outputs valid JSON.
Do not add explanations, markdown, or text outside the JSON.
Output format example:
[
  {"title": "Slide Title", "bullets": ["point 1", "point 2"]},
  {"title": "Another Slide", "bullets": ["bullet a", "bullet b"]}
]""" },
                    {"role": "user", "content": f"Input text:\n{text_input}\n\nGuidance:{guidance}\n\nReturn only the JSON slide plan."}
                ]

                resp = generate_completion_with_provider(provider, api_key, model="gpt-4o-mini", messages=messages)

                # Different providers nest responses differently
                if provider in ["openai", "aipipe", "aipipe_openrouter"]:
                    content = resp["choices"][0]["message"]["content"]
                elif provider == "anthropic":
                    content = resp["content"][0]["text"]
                elif provider == "gemini":
                    content = resp["candidates"][0]["content"]["parts"][0]["text"]
                else:
                    content = "[]"

                # Debug: show raw response content
                st.text_area("DEBUG LLM Output", content, height=200)

                # Parse JSON safely
                slide_plan = safe_json_parse(content, [{"title": "Summary", "bullets": [text_input[:200] + "..."]}])

                prs = build_presentation(uploaded_template, slide_plan)
                output = io.BytesIO()
                prs.save(output)
                st.success("✅ Presentation ready!")
                st.download_button("Download PPTX", data=output.getvalue(), file_name="generated.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

            finally:
                # clear key
                del api_key
